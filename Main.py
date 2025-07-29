import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.dml.color import RGBColor
import re
from dataclasses import dataclass
from typing import List, Dict, Any, Tuple
import os
import win32com.client
import time
import math
import copy

# =================================================================================
# SHARED HELPER FUNCTIONS
# =================================================================================

def find_shape_by_name(slide, name):
    """Finds a shape on a slide by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"Warning: Shape with name '{name}' not found on the slide.")
    return None

def format_main_title(shape, last_week):
    """Sets the text and formatting for the main title of the slide."""
    if not (shape and shape.has_text_frame): return
    match = re.search(r'Week\s*(\d+)', str(last_week))
    week_num = match.group(1) if match else ""
    if not week_num:
        print("Warning: Could not extract week number from title.")
        return
    
    original_text = shape.text_frame.text
    
    if '$' in original_text:
        new_text = original_text.replace('$', week_num)
    else:
        new_text = re.sub(r"(past)(\s+)(weeks)", rf"\1 {week_num} \3", original_text, flags=re.IGNORECASE)

    p = shape.text_frame.paragraphs[0]
    p.text = new_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True

# =================================================================================
# LOGIC FOR SLIDE 6
# =================================================================================

@dataclass
class TableData:
    title: str
    headers: List[str]
    row_labels: List[str]
    data: Dict[str, List[Any]]

def _extract_slide6_table_data(df: pd.DataFrame, title: str, header_row_idx: int, num_data_rows: int) -> TableData:
    header_row = df.iloc[header_row_idx]
    all_weeks = [header_row[col] for col in df.columns if "Week" in str(header_row[col])]
    last_four_weeks = all_weeks[-4:]
    week_cols = [col for col in df.columns if header_row[col] in last_four_weeks]
    data_start_row = header_row_idx + 1
    data_end_row = data_start_row + num_data_rows
    data_rows_df = df.iloc[data_start_row:data_end_row].set_index(df.columns[0])
    row_labels = data_rows_df.index.tolist()
    data_dict = {header_row[col]: data_rows_df[col].tolist() for col in week_cols}
    return TableData(title=title, headers=last_four_weeks, row_labels=row_labels, data=data_dict)

def _set_font_for_table_slide6(table):
    for row_idx, row in enumerate(table.rows):
        is_total_row = 'total' in str(row.cells[0].text).lower()
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(11)
                    if row_idx == 0 or is_total_row:
                        run.font.bold = True
                    else:
                        run.font.bold = False

def _populate_table_slide6(table_shape, table_data: TableData):
    if not (table_shape and table_shape.has_table): return
    table = table_shape.table
    for i, header in enumerate(table_data.headers):
        table.cell(0, i + 1).text = str(header)
    for row_idx, label in enumerate(table_data.row_labels):
        table.cell(row_idx + 1, 0).text = label
        for col_idx, header in enumerate(table_data.headers):
            table.cell(row_idx + 1, col_idx + 1).text = str(table_data.data[header][row_idx])
    _set_font_for_table_slide6(table)

def _add_line_chart_slide6(slide, position, table_data: TableData, num_gridlines: int):
    chart_data = CategoryChartData()
    chart_data.categories = table_data.headers
    
    all_values = []
    for i, label in enumerate(table_data.row_labels):
        if 'total' in label.lower(): continue
        values = [table_data.data[h][i] for h in table_data.headers]
        all_values.extend(values)
        chart_data.add_series(label, values)
        
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, *position, chart_data)
    chart = chart_shape.chart
    
    value_axis = chart.value_axis
    if num_gridlines and all_values:
        max_val = max(all_values) if all_values else 0
        if max_val > 0:
            power = 10**math.floor(math.log10(max_val))
            max_y = math.ceil(max_val / power) * power
            if max_y % num_gridlines != 0:
                max_y = math.ceil(max_y / num_gridlines) * num_gridlines
        else:
            max_y = float(num_gridlines * 2)

        value_axis.maximum_scale = max_y
        value_axis.minimum_scale = 0.0
        value_axis.major_unit = max_y / num_gridlines
        value_axis.has_major_gridlines = True

    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0'
    tick_labels.font.size = Pt(9)
    
    chart.has_title = True
    chart.chart_title.text_frame.text = f'{table_data.title} Resolved'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(22)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)

def populate_slide_6(prs, df, slide_index):
    print(f"--- Populating Slide {slide_index + 1} ---")
    slide = prs.slides[slide_index]
    inc_data_obj = _extract_slide6_table_data(df, title="INC", header_row_idx=0, num_data_rows=5)
    ritm_data_obj = _extract_slide6_table_data(df, title="RITM", header_row_idx=8, num_data_rows=5)
    print("   -> Formatting main title...")
    title_shape = find_shape_by_name(slide, 'MainTitle')
    if title_shape: format_main_title(title_shape, inc_data_obj.headers[-1])
    print("   -> Populating tables...")
    inc_table_shape = find_shape_by_name(slide, 'INC Table') 
    ritm_table_shape = find_shape_by_name(slide, 'RITM Table')
    _populate_table_slide6(inc_table_shape, inc_data_obj)
    _populate_table_slide6(ritm_table_shape, ritm_data_obj)
    for shape in list(slide.shapes):
        if shape.has_chart:
            sp = shape._sp
            sp.getparent().remove(sp)

    inc_chart_position = (Cm(0.67), Cm(1.81), Cm(15.49), Cm(6.58))
    ritm_chart_position = (Cm(17.42), Cm(1.81), Cm(15.49), Cm(6.58))
    
    print("   -> Creating line charts...")
    _add_line_chart_slide6(slide, inc_chart_position, inc_data_obj, num_gridlines=5)
    _add_line_chart_slide6(slide, ritm_chart_position, ritm_data_obj, num_gridlines=5)
    print(f"✅ Slide {slide_index + 1} populated.")

# =================================================================================
# LOGIC FOR SLIDE 7
# =================================================================================

def _find_data_row(df: pd.DataFrame, label_text: str) -> int:
    for index, row in df.iterrows():
        if label_text.lower() in str(row.iloc[0]).lower():
            return index
    return -1

def _extract_data_block_slide7(df: pd.DataFrame, start_label: str, num_rows: int):
    start_row_idx = _find_data_row(df, start_label)
    if start_row_idx == -1: return None, None
    header_row_idx = start_row_idx - 1
    header_row = df.iloc[header_row_idx]
    all_week_cols = [col for col, h in header_row.items() if "Week" in str(h)]
    last_four_week_cols = all_week_cols[-4:]
    last_four_week_headers = [str(header_row[col]) for col in last_four_week_cols]
    data_end_row = start_row_idx + num_rows
    categories = df.iloc[start_row_idx : data_end_row, 0].values.tolist()
    data_dict = {h: df.loc[start_row_idx : data_end_row - 1, col].tolist() for h, col in zip(last_four_week_headers, last_four_week_cols)}
    return categories, data_dict

def _extract_stats_data_slide7(df: pd.DataFrame):
    inc_header_row = df.iloc[0]
    all_week_cols = [col for col, h in inc_header_row.items() if "Week" in str(h)]
    last_four_week_cols = all_week_cols[-4:]
    last_four_week_headers = [str(inc_header_row[col]) for col in last_four_week_cols]
    created_row_idx = _find_data_row(df, 'INCs created')
    resolved_row_idx = _find_data_row(df, 'INCs resolved')
    if created_row_idx == -1 or resolved_row_idx == -1: return None, None
    created_data = df.loc[created_row_idx, last_four_week_cols].tolist()
    resolved_data = df.loc[resolved_row_idx, last_four_week_cols].tolist()
    categories = ['INCs created', 'INCs resolved']
    data_dict = {w: [c, r] for w, c, r in zip(last_four_week_headers, created_data, resolved_data)}
    return categories, data_dict

def _add_bar_chart_slide7(slide, categories, data_dict, position, title, y_axis_max=None, num_gridlines=None):
    if data_dict is None: return
    chart_data = ChartData()
    chart_data.categories = categories
    for week, values in data_dict.items():
        chart_data.add_series(week, values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *position, chart_data)
    chart = chart_shape.chart
    
    value_axis = chart.value_axis
    if y_axis_max and num_gridlines:
        value_axis.maximum_scale = float(y_axis_max)
        value_axis.major_unit = float(y_axis_max) / num_gridlines
    elif y_axis_max:
        value_axis.maximum_scale = float(y_axis_max)

    value_axis.has_major_gridlines = True
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0'
    tick_labels.font.size = Pt(9)

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)

def populate_slide_7(prs, df, slide_index):
    print(f"--- Populating Slide {slide_index + 1} ---")
    slide = prs.slides[slide_index]
    print("   -> Extracting data for bar charts...")
    creation_cats, creation_data = _extract_data_block_slide7(df, start_label="Tools Created", num_rows=4)
    closed_cats, closed_data = _extract_data_block_slide7(df, start_label="Auto closed by Tools", num_rows=2)
    stats_cats, stats_data = _extract_stats_data_slide7(df)
    pos_created = (Cm(1.69), Cm(1.76), Cm(13.27), Cm(6.38))
    pos_resolved = (Cm(18.27), Cm(1.76), Cm(13.91), Cm(6.38))
    pos_stats = (Cm(1.69), Cm(9.75), Cm(13.27), Cm(6.77))
    print("   -> Creating bar charts with fixed axis scale...")
    _add_bar_chart_slide7(slide, creation_cats, creation_data, pos_created, "INC Created by", y_axis_max=1600.0, num_gridlines=7)
    _add_bar_chart_slide7(slide, closed_cats, closed_data, pos_resolved, "INC Resolved by", y_axis_max=1600.0, num_gridlines=7)
    _add_bar_chart_slide7(slide, stats_cats, stats_data, pos_stats, "INC Stats", y_axis_max=1600.0, num_gridlines=7)
    print(f"✅ Slide {slide_index + 1} populated.")

# =================================================================================
# LOGIC FOR SLIDE 8
# =================================================================================

def _extract_pending_data_slide8(df: pd.DataFrame, series_labels: List[str], week_column_label: str = 'Week #') -> Tuple[List[str], Dict[str, List]]:
    valid_series_labels = [label for label in series_labels if label in df.columns]
    if not valid_series_labels: return None, None
    last_four_rows_df = df.tail(4)
    categories = last_four_rows_df[week_column_label].astype(str).tolist()
    data_dict = {label: last_four_rows_df[label].tolist() for label in valid_series_labels}
    return categories, data_dict

def _add_chart_slide8(slide, categories, data_dict, position, title, num_gridlines: int):
    if data_dict is None: return
    chart_data = ChartData()
    chart_data.categories = categories
    
    all_values = []
    for series_name, values in data_dict.items():
        chart_data.add_series(series_name, values)
        all_values.extend(values)

    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, *position, chart_data)
    chart = chart_shape.chart
    
    value_axis = chart.value_axis
    if num_gridlines and all_values:
        max_val = max(all_values) if all_values else 0
        if max_val > 0:
            power = 10**math.floor(math.log10(max_val))
            max_y = math.ceil(max_val / power) * power
            if max_y % num_gridlines != 0:
                max_y = math.ceil(max_y / num_gridlines) * num_gridlines
        else:
            max_y = float(num_gridlines * 10)

        value_axis.maximum_scale = max_y
        value_axis.minimum_scale = 0.0
        value_axis.major_unit = max_y / num_gridlines
        value_axis.has_major_gridlines = True

    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0'
    tick_labels.font.size = Pt(9)
    
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(22)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.category_axis.tick_labels.font.size = Pt(9)

def _update_text_boxes_slide8(slide, inc_data, ritm_data, inc_cats, ritm_cats):
    if inc_data:
        inc_total_created = inc_data.get("Total INCs Created", [0])[-1]
        last_week_num = re.search(r'\d+', inc_cats[-1]).group(0)
        inc_conc_shape = find_shape_by_name(slide, "INC Conc")
        if inc_conc_shape and inc_conc_shape.has_text_frame:
            p1 = inc_conc_shape.text_frame.paragraphs[0]
            p1.text = re.sub(
                r'(Total Incidents raised for the week\s+)\d+(\s+is\s+)\d+(\.)',
                rf'\g<1>{last_week_num}\g<2>{inc_total_created}\g<3>',
                p1.text
            )
            for run in p1.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(11)
            print("   -> Updated 'INC Conc' text box.")

    if ritm_data:
        ritm_total_created = ritm_data.get("Total RITMs Created", [0])[-1]
        ritm_fulfilled = ritm_data.get("RITMs Fulfilled", [0])[-1]
        last_week_num = re.search(r'\d+', ritm_cats[-1]).group(0)
        ritm_conc_shape = find_shape_by_name(slide, "RITM Conc")
        if ritm_conc_shape and ritm_conc_shape.has_text_frame:
            p = ritm_conc_shape.text_frame.paragraphs[0]
            p.text = re.sub(
                r'(Total RITM raised for the Week\s+)\d+(\s+is\s+)\d+(,\s*Fulfilled\s*–\s*)\d+',
                rf'\g<1>{last_week_num}\g<2>{ritm_total_created}\g<3>{ritm_fulfilled}',
                p.text
            )
            for run in p.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(11)
            print("   -> Updated 'RITM Conc' text box.")

def populate_slide_8(prs, df, slide_index):
    print(f"--- Populating Slide {slide_index + 1} ---")
    slide = prs.slides[slide_index]
    df.columns = df.columns.str.strip()
    print("   -> Extracting data for charts...")
    inc_series_labels = ["Pending INCs", "INCs Resolved", "Total INCs Created"]
    ritm_series_labels = ["Pending RITMs", "RITMs Fulfilled", "Total RITMs Created"]
    inc_cats, inc_data = _extract_pending_data_slide8(df, inc_series_labels)
    ritm_cats, ritm_data = _extract_pending_data_slide8(df, ritm_series_labels)
    
    # Updated positions for Slide 8 charts
    pos_inc = (Cm(1.94), Cm(2.28), Cm(13.65), Cm(8.1))
    pos_ritm = (Cm(18.01), Cm(2.28), Cm(13.65), Cm(8.1))

    print("   -> Creating line charts...")
    _add_chart_slide8(slide, inc_cats, inc_data, pos_inc, "INC Resolved", num_gridlines=8)
    _add_chart_slide8(slide, ritm_cats, ritm_data, pos_ritm, "RITMs Resolved", num_gridlines=8)
    print("   -> Updating text boxes...")
    _update_text_boxes_slide8(slide, inc_data, ritm_data, inc_cats, ritm_cats)
    print(f"✅ Slide {slide_index + 1} populated.")

# =================================================================================
# LOGIC FOR SLIDE 9
# =================================================================================

def _extract_data_block_slide9(df: pd.DataFrame, start_label: str, num_rows: int):
    start_row_idx = _find_data_row(df, start_label)
    if start_row_idx == -1: return None, None, None
    header_row = df.iloc[start_row_idx]
    all_data_cols = df.columns[1:]
    last_four_data_cols = all_data_cols[-4:]
    last_four_week_headers = [str(header_row[col]) for col in last_four_data_cols]
    data_start_row = start_row_idx + 1
    data_end_row = data_start_row + num_rows
    categories = df.iloc[data_start_row : data_end_row, 0].values.tolist()
    data_dict = {h: df.loc[data_start_row : data_end_row - 1, col].tolist() for h, col in zip(last_four_week_headers, last_four_data_cols)}
    return categories, data_dict, last_four_week_headers

def _set_font_for_table_slide9(table):
    for row_idx, row in enumerate(table.rows):
        is_total_row = 'total' in str(row.cells[0].text).lower()
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(11)
                    if row_idx == 0 or is_total_row:
                        run.font.bold = True
                    else:
                        run.font.bold = False

def _populate_table_slide9(table_shape, categories, data_dict, headers, top_left_cell_text: str):
    if not (table_shape and table_shape.has_table): return
    table = table_shape.table
    table.cell(0, 0).text = top_left_cell_text
    for i, header in enumerate(headers):
        table.cell(0, i + 1).text = str(header)
    for row_idx, label in enumerate(categories):
        table.cell(row_idx + 1, 0).text = label
        for col_idx, header in enumerate(headers):
            table.cell(row_idx + 1, col_idx + 1).text = str(data_dict[header][row_idx])
    _set_font_for_table_slide9(table)

def _add_line_chart_slide9(slide, categories, data_dict, position, title, num_gridlines: int):
    if data_dict is None: return
    chart_data = ChartData()
    chart_data.categories = [c for c in categories if 'total' not in c.lower()]
    
    all_values = []
    for week, values in data_dict.items():
        series_values = [val for cat, val in zip(categories, values) if 'total' not in cat.lower()]
        all_values.extend(series_values)
        chart_data.add_series(week, series_values)

    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, *position, chart_data)
    chart = chart_shape.chart
    
    value_axis = chart.value_axis
    if num_gridlines and all_values:
        max_val = max(all_values) if all_values else 0
        if max_val > 0:
            power = 10**math.floor(math.log10(max_val))
            max_y = math.ceil(max_val / power) * power
            if max_y % num_gridlines != 0:
                max_y = math.ceil(max_y / num_gridlines) * num_gridlines
        else:
            max_y = float(num_gridlines * 5)

        value_axis.maximum_scale = max_y
        value_axis.minimum_scale = 0.0
        value_axis.major_unit = max_y / num_gridlines
        value_axis.has_major_gridlines = True

    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0'
    tick_labels.font.size = Pt(9)

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(22)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)

def _update_conclusion_slide9(slide, categories, data_dict, headers):
    if data_dict:
        total_row_index = [i for i, cat in enumerate(categories) if 'total' in cat.lower()]
        if not total_row_index: return
        total_row_index = total_row_index[0]
        
        last_week_header = headers[-1]
        last_week_total = data_dict[last_week_header][total_row_index]
        last_week_num = re.search(r'\d+', last_week_header).group(0)

        conclusion_shape = find_shape_by_name(slide, "Conclusion")
        if conclusion_shape and conclusion_shape.has_text_frame:
            p = conclusion_shape.text_frame.paragraphs[0]
            p.text = f"{last_week_total} successful changes for Week {last_week_num}"
            print("   -> Updated 'Conclusion' text box.")

def populate_slide_9(prs, df, slide_index):
    print(f"--- Populating Slide {slide_index + 1} ---")
    slide = prs.slides[slide_index]
    
    print("   -> Extracting data for tables and charts...")
    type_cats, type_data, type_headers = _extract_data_block_slide9(df, start_label="Type", num_rows=6)
    closure_cats, closure_data, closure_headers = _extract_data_block_slide9(df, start_label="Closure Type", num_rows=4)

    print("   -> Formatting main title...")
    title_shape = find_shape_by_name(slide, 'Title')
    if title_shape and type_headers:
        format_main_title(title_shape, type_headers[-1])

    print("   -> Populating tables...")
    cr_table_shape = find_shape_by_name(slide, "CR Table")
    ccs_table_shape = find_shape_by_name(slide, "CCS Table")
    _populate_table_slide9(cr_table_shape, type_cats, type_data, type_headers, "Change")
    _populate_table_slide9(ccs_table_shape, closure_cats, closure_data, closure_headers, "Closure")

    pos_cr = (Cm(0.78), Cm(1.64), Cm(15.31), Cm(7.63))
    pos_ccs = (Cm(16.54), Cm(1.64), Cm(16.55), Cm(7.6))
    
    print("   -> Creating line charts...")
    _add_line_chart_slide9(slide, type_cats, type_data, pos_cr, "Changes Requests", num_gridlines=7)
    _add_line_chart_slide9(slide, closure_cats, closure_data, pos_ccs, "Change Closure Status", num_gridlines=7)

    print("   -> Updating conclusion text box...")
    _update_conclusion_slide9(slide, closure_cats, closure_data, closure_headers)

    print(f"✅ Slide {slide_index + 1} populated.")

# =================================================================================
# LOGIC FOR SLIDE 10
# =================================================================================

def _extract_data_for_slide10(df: pd.DataFrame):
    """Extracts all table data from the 'Business Services' sheet."""
    all_tables_data = []
    target_table_titles = ["Business Service - INC", "Business Service - RITM", "Business Service - Change Requests"]

    for table_title in target_table_titles:
        try:
            df[0] = df[0].astype(str).str.strip()
            title_rows = df[df[0] == table_title]
            if title_rows.empty: continue
            start_index = title_rows.index[0]

            search_area = df.iloc[start_index:]
            grand_total_rows = search_area[search_area[0] == 'Grand Total']
            if grand_total_rows.empty: continue
            end_index = grand_total_rows.index[0]

            header_row_index = -1
            for idx in range(start_index, end_index):
                if any('Week' in str(cell) for cell in df.iloc[idx].iloc[-4:]):
                    header_row_index = idx
                    break
            if header_row_index == -1: continue
            
            header_row = df.iloc[header_row_index]
            week_labels = [str(label).strip() for label in header_row.iloc[-4:]]
            
            headers = [table_title] + week_labels
            data_rows = []
            for i in range(header_row_index + 1, end_index + 1):
                current_row = df.iloc[i]
                row_title = str(current_row.iloc[0]).strip()
                if not row_title: continue
                data = [int(val) for val in current_row.iloc[-4:]]
                data_rows.append([row_title] + data)
            
            all_tables_data.append({"title": table_title, "headers": headers, "data": data_rows})
        except Exception as e:
            print(f"  -> ERROR processing table '{table_title}': {e}")
    return all_tables_data

def _create_and_populate_table_slide10(slide, position: dict, table_data: dict):
    """Creates, populates, and formats a new table to fit exactly within the given dimensions."""
    headers = table_data["headers"]
    data_rows = table_data["data"]
    num_rows, num_cols = len(data_rows) + 1, len(headers)
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, **position)
    table = table_shape.table

    table.columns[0].width = int(position['width'] * 0.60)
    data_col_width = int((position['width'] * 0.40) / (num_cols - 1))
    for i in range(1, num_cols):
        table.columns[i].width = data_col_width

    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header

    for row_idx, data_row in enumerate(data_rows, 1):
        for col_idx, cell_text in enumerate(data_row):
            table.cell(row_idx, col_idx).text = str(cell_text)

    for row_idx, row in enumerate(table.rows):
        is_total_row = 'total' in table.cell(row_idx, 0).text.lower()
        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.1)
            cell.margin_bottom = Cm(0.1)

            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(7)
                    if row_idx == 0 or is_total_row:
                        run.font.bold = True

def populate_slide_10(prs, df, slide_index):
    """Main function to create and populate the three tables on Slide 10."""
    print(f"--- Populating Slide {slide_index + 1} ---")
    slide = prs.slides[slide_index]
    
    all_tables_data = _extract_data_for_slide10(df)
    if not all_tables_data:
        print(f"🛑 Could not extract data for Slide {slide_index + 1}. Skipping.")
        return

    table_positions = {
        "Business Service - INC": {"left": Cm(0.78), "top": Cm(1.89), "width": Cm(16.07), "height": Cm(8.22)},
        "Business Service - RITM": {"left": Cm(17.35), "top": Cm(1.89), "width": Cm(15.73), "height": Cm(8.51)},
        "Business Service - Change Requests": {"left": Cm(0.78), "top": Cm(10.55), "width": Cm(16.07), "height": Cm(6.95)}
    }
    
    for shape_name in ["BSI", "BSR", "BSC"]:
        if shape := find_shape_by_name(slide, shape_name):
            sp = shape._sp
            sp.getparent().remove(sp)
            print(f"  -> Deleted placeholder table '{shape_name}'.")

    for table_data in all_tables_data:
        if position := table_positions.get(table_data["title"]):
            print(f"  -> Creating table for '{table_data['title']}'...")
            _create_and_populate_table_slide10(slide, position, table_data)
        else:
            print(f"  -> WARNING: No position mapping for '{table_data['title']}'.")

    print(f"✅ Slide {slide_index + 1} populated.")

# =================================================================================
# REUSABLE STYLING FUNCTION
# =================================================================================

def apply_chart_styles(input_path, output_path, styles_to_apply: Dict[int, int]):
    print(f"\n--- Applying Final Chart Styles ---")
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
        powerpoint.WindowState = 2
        time.sleep(2)
        for slide_index, style_id in styles_to_apply.items():
            if presentation.Slides.Count < (slide_index + 1):
                print(f"   -> Skipping style for slide {slide_index + 1}, slide does not exist.")
                continue
            target_slide = presentation.Slides(slide_index + 1)
            print(f"   -> Targeting Slide {slide_index + 1} for Style ID {style_id}...")
            for shape in target_slide.Shapes:
                if shape.HasChart:
                    shape.Chart.ChartStyle = style_id
        presentation.SaveAs(os.path.abspath(output_path))
    finally:
        if presentation: presentation.Close()
        if powerpoint: powerpoint.Quit()
    print("✅ Styling complete.")

# =================================================================================
# MAIN EXECUTION WORKFLOW
# =================================================================================

if __name__ == "__main__":
    EXCEL_FILE_PATH = 'Files/ExcelData27.xlsx'
    TEMPLATE_PPTX_PATH = 'Files/Default_Template.pptx'
    FINAL_OUTPUT_PPTX_PATH = 'Output/Final_Output.pptx'

    SHEET_NAME_SLIDE_6 = 'Volumetric trends INC & RITM'
    SHEET_NAME_SLIDE_7 = 'Created'
    SHEET_NAME_SLIDE_8 = 'Pending Counts'
    SHEET_NAME_SLIDE_9 = 'Volumetric Change details'
    SHEET_NAME_SLIDE_10 = 'Business Services'
    
    SLIDE_6_INDEX = 2
    SLIDE_7_INDEX = 3
    SLIDE_8_INDEX = 4
    SLIDE_9_INDEX = 5
    SLIDE_10_INDEX = 6

    CHART_STYLE_FOR_LINE_CHARTS = 228

    os.makedirs(os.path.dirname(FINAL_OUTPUT_PPTX_PATH), exist_ok=True)
    temp_output_path = os.path.join(os.path.dirname(FINAL_OUTPUT_PPTX_PATH), "temp_presentation.pptx")

    try:
        df_slide6 = pd.read_excel(EXCEL_FILE_PATH, sheet_name=SHEET_NAME_SLIDE_6)
        df_slide7 = pd.read_excel(EXCEL_FILE_PATH, sheet_name=SHEET_NAME_SLIDE_7)
        df_slide8 = pd.read_excel(EXCEL_FILE_PATH, sheet_name=SHEET_NAME_SLIDE_8)
        df_slide9 = pd.read_excel(EXCEL_FILE_PATH, sheet_name=SHEET_NAME_SLIDE_9)
        df_slide10 = pd.read_excel(EXCEL_FILE_PATH, sheet_name=SHEET_NAME_SLIDE_10, header=None).fillna(0)
    except Exception as e:
        print(f"🛑 FATAL ERROR reading Excel file: {e}")
        exit()

    prs = Presentation(TEMPLATE_PPTX_PATH)
    
    populate_slide_6(prs, df_slide6, SLIDE_6_INDEX)
    populate_slide_7(prs, df_slide7, SLIDE_7_INDEX)
    populate_slide_8(prs, df_slide8, SLIDE_8_INDEX)
    populate_slide_9(prs, df_slide9, SLIDE_9_INDEX)
    populate_slide_10(prs, df_slide10, SLIDE_10_INDEX)
    
    print("\n💾 Saving presentation with all content...")
    prs.save(temp_output_path)
    
    styles_to_apply = {
        SLIDE_6_INDEX: CHART_STYLE_FOR_LINE_CHARTS,
        SLIDE_8_INDEX: CHART_STYLE_FOR_LINE_CHARTS,
        SLIDE_9_INDEX: CHART_STYLE_FOR_LINE_CHARTS
    }

    apply_chart_styles(
        input_path=temp_output_path,
        output_path=FINAL_OUTPUT_PPTX_PATH,
        styles_to_apply=styles_to_apply
    )
    
    try:
        os.remove(temp_output_path)
        print(f"🧹 Temporary file has been deleted.")
    except OSError as e:
        print(f"Warning: Could not delete temporary file: {e}")

    print(f"\n🎉🎉🎉 Workflow Finished! Your final presentation is ready at:\n{FINAL_OUTPUT_PPTX_PATH}")
